import os
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from research_parser import research_context_parser, research_figure_parser
# from conference_parser import conference_parser
from models import PosterBulletPoints
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate

# poster formatter agent
def poster_formatter_agent():
    """
    Agent that generates bullet points for each section for poster
    """
    
    llm = ChatOpenAI(
        model="gpt-4o-mini",
        temperature= 0
    )
    structured_llm = llm.with_structured_output(PosterBulletPoints)
    
    system_instructions = """
    You are an expert Presentation Designer.
    Your job is to take academic paragraphs or research abstract and condense them into extremely concise, short, punchy bullet points suitable for a massive 48x36 scientific poster.
    
    Rules:
    1. Do not lose key metrics or numbers.
    2. Remove filler words (e.g. change "We utilized a novel approach to" to "Utilized a novel approach for" )
    3. Ensure maximum scannability.
    """
    prompt = ChatPromptTemplate([
        ("system", system_instructions),
        ("user", """Format this research data for a poster:\n Title: {title}\nIntro: {intro}\nProblem: {problem}
         \nResearch Goals: {research_goals}\nMethods: {methods}\nResults: {results}\n Conclusion: {conclusion}
         """)
    ])
    return prompt | structured_llm
    
def add_section(slide, title, bullet_list, left, top, width, body_size, title_size):
    """
    Helper func to draw a text box with a bold header

    Args:
        slide : the slide object
        title (str): title of the research section
        bulet_list(list): A list of strings to be rendered as bullet points
        left, top, width: Positions in Inches
        body_size, title_size: Font sizes
    """
    if not bullet_list:
        return top
    
    # add text box to the poster
    text_box = slide.shapes.add_textbox(left, top, width, Inches(1))
    
    # add text to the text box frame
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # Header 
    p_heading = text_frame.paragraphs[0]
    p_heading.text = title
    p_heading.font.bold = True
    p_heading.font.size = Pt(title_size)
    
    # body - loop thru the bullet points
    for bullet_text in bullet_list:
        p_body = text_frame.add_paragraph()
        p_body.text = bullet_text
        p_body.font.size = Pt(body_size)
        
        p_body.level = 1 # turns the paragraph into a bullet point!
    
    # Estimate height to push the next section down
    # 1. Sum up all the chars in the list
    total_chars = sum(len(b) for b in bullet_list)
    # 2. Calculate height: base header height + (text wrapping math) + a little padding per bullet
    estimated_hei = Inches(1.5) + (total_chars / 50 * Inches(0.3)) + (len(bullet_list) * Inches(0.1))
    
    return top + estimated_hei

def main():
    sample_raw_research_text = """
    Date: October 14
    Phase: Initial Setup & Baseline Testing

    Summary of Progress:
    During this phase, we established the baseline Retrieval-Augmented Generation (RAG) pipeline to assist students with queries related to data structures and algorithms. The system ingests university lecture slides, textbook chapters, and documentation. We utilized a standard vector database for semantic search and a pre-trained Large Language Model (LLM) to generate answers based on the retrieved context.
    Initial Findings:
    We evaluated the baseline system on a test set of 500 algorithmic questions.

    Overall Accuracy: 68.5%

    Observation: The model performs well on direct factual recall (e.g., "What is the time complexity of QuickSort?"). However, it struggles significantly with multi-hop reasoning and constraint-based queries (e.g., "Which sorting algorithm is best for a nearly sorted array with strict memory constraints?"). Hallucinations occur in 18% of complex logic responses.

    Date: November 4
    Phase: Architecture Revision & Neuro-Symbolic Integration

    Summary of Progress:
    To address the LLM's failure in complex reasoning, we redesigned the pipeline into a neuro-symbolic architecture. Instead of relying solely on the LLM to generate the final answer from text, the LLM now acts as a semantic parser. It extracts the rules and constraints from the retrieved RAG context and formats them into logical predicates. A deterministic symbolic solver then evaluates these rules to guarantee a logically sound answer.

    Overall Accuracy: Improved to 84.2%.

    Observation: Hallucinations on complex algorithmic questions dropped to 3%. The symbolic solver successfully enforces strict logical consistency.

    Issue: The system latency increased significantly. Processing queries through the symbolic solver adds an average of 850ms per query, making the system sluggish for simple questions that do not require deep reasoning.

    Date: November 25
    Phase: Optimization & Final Benchmarking

    Summary of Progress:
    To solve the latency bottleneck identified in the previous phase, we implemented a lightweight query classification routing mechanism. The system now uses a small, fine-tuned classifier to determine if a student's query is "factual" or "logical/algorithmic." Factual queries bypass the symbolic engine and use the standard baseline RAG, while complex queries are routed to the neuro-symbolic pipeline.

    Final Accuracy: 88.7% across the entire 500-question test set.

    Latency: Average response time reduced to 320ms, which is well within the acceptable threshold for real-time educational tools.

    Conclusion: The hybrid routing approach successfully maintains the high accuracy and logical consistency of neuro-symbolic AI while preserving the speed and efficiency of traditional RAG systems.

    """

    # A traditional biopsy method used for investigating cancer-suspicious skin tissue, while effective, is invasive and time-consuming. It includes multiple processing steps to produce Hematoxylin & Eosin (H&E) images, which are gold standard for cancer diagnosis. Recently, a virtual biopsy has emerged as a non-invasive alternative that produces two-dimensional grayscale Optical Coherence Tomography (OCT) images by simply scanning live tissue. However, interpreting OCT images requires extensive specialized training for pathologists, which makes the virtual biopsy limited in clinical settings. Moreover, the scarcity of annotated OCT datasets has hindered the development of useful foundation models to assist pathologists examining these OCT images. In this study, we developed a vision-language model capable of interpreting OCT images for tissue classifications without manual annotations. We fine-tuned the Multimodal transformer with Unified maSKed modeling (MUSK), a vision-language model originally pre-trained on H&E images, using a knowledge distillation approach. By training the model on approximately 500 H&E-OCT image pairs, we reduced the domain gap between H&E image embeddings (feature representations) from original MUSK vision encoder (teacher model) and OCT image embeddings from fine-tuned MUSK vision encoder (student model). We achieved around 72% accuracy increase in skin tissue classification in OCT images. This significant gain in accuracy highlights the model’s potential to support pathologists in interpreting OCT scans and enhance the clinical viability of non-invasive cancer diagnosis.


    user_focus = "poster"
    with open("./RAG.png", "rb") as image_file:
        encoded_str = base64.b64encode(image_file.read()).decode('utf-8')
        # add data URI prefix to tell api what kind of media files those chars represent;; otherwise the api will reject it or read it as gibberish
        formatted_image_uri = f"data:image/png;base64,{encoded_str}"
        
    with open("./RAG_2.png", "rb") as image_file:
        encoded_str2 = base64.b64encode(image_file.read()).decode('utf-8')
        # add data URI prefix to tell api what kind of media files those chars represent;; otherwise the api will reject it or read it as gibberish
        formatted_image_uri2 = f"data:image/png;base64,{encoded_str2}"
        
    with open("./RAG_3.png", "rb") as image_file:
        encoded_str3 = base64.b64encode(image_file.read()).decode('utf-8')
        # add data URI prefix to tell api what kind of media files those chars represent;; otherwise the api will reject it or read it as gibberish
        formatted_image_uri3 = f"data:image/png;base64,{encoded_str3}"
        
    extracted_images = {
        'RAG': formatted_image_uri, 
        'RAG_2': formatted_image_uri2,
        'RAG_3': formatted_image_uri3
        }
    # extracted_images = None
    if extracted_images:
        vision_metadata = research_figure_parser(sample_raw_research_text, extracted_images)
    else:
        vision_metadata = None
    # print(vision_metadata)

    for img in vision_metadata:
        print(f"File: {img.image_filename}")
        print(f"Section: {img.suggested_section}")
        print(f"Caption: {img.figure_caption}")
        print(f"Deep Analysis: {img.detailed_analysis}")
        print()

    agent = research_context_parser()
    # run the agent 
    print("running research parser...")
    res = agent.invoke({
        "user_goal": user_focus,
        "raw_text": sample_raw_research_text,
        "vision_metadata": vision_metadata
    })

    print("\n--- Parsed Output ---")
    print(f"Title: {res.title}\n")
    print(f"Introduction: {res.introduction}")
    print(f"Problem Gap: {res.problem_gap}\n")
    print(f"Research Goal: {res.research_goal}\n")
    print(f"Methodologies: {res.methodology}\n")
    print(f"Results: {res.key_res}\n")
    print(f"Conclusion: {res.conclusion}\n")
    print(f"Abstract: {res.abstract}\n")
    print()

    poster_agent = poster_formatter_agent()
    bullet_points = poster_agent.invoke({
        'title': res.title,
        'intro': res.introduction,
        'problem': res.problem_gap,
        'research_goals': res.research_goal,
        'methods': res.methodology,
        'results': res.key_res,
        'conclusion': res.conclusion
    })

    print("Parsiing bullet points for poster")
    print("\n--- Parsed Bullet Points for poster ---")
    print(f"Title: {bullet_points.title}\n")
    print(f"Introduction: {bullet_points.intro_bullets}")
    print(f"Problem Gap: {bullet_points.problem_gap_bullets}\n")
    print(f"Research Goal: {bullet_points.research_goal_bullets}\n")
    print(f"Methodologies: {bullet_points.method_bullets}\n")
    print(f"Results: {bullet_points.result_bullets}\n")
    print(f"Conclusion: {bullet_points.conclusion_bullets}\n")
    print()
        
if __name__ == '__main__':
    main()